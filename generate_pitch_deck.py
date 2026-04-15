#!/usr/bin/env python3
"""
Pacífico Bambu — Investor Pitch Deck Generator
Generates a professional 16-slide pitch deck using python-pptx.
Dark green theme matching reference design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme Colors ──────────────────────────────────────────────
BG_DARK    = RGBColor(0x0A, 0x2E, 0x1A)   # Dark green background
BG_CARD    = RGBColor(0x12, 0x3D, 0x26)   # Card background
ACCENT     = RGBColor(0x2D, 0xD4, 0xA8)   # Mint green accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_80   = RGBColor(0xCC, 0xCC, 0xCC)
WHITE_60   = RGBColor(0x99, 0x99, 0x99)
WHITE_40   = RGBColor(0x66, 0x66, 0x66)
WHITE_20   = RGBColor(0x33, 0x33, 0x33)
RED_SOFT   = RGBColor(0xF0, 0x6E, 0x6E)
YELLOW     = RGBColor(0xFA, 0xCC, 0x15)
BLUE_SOFT  = RGBColor(0x60, 0xA5, 0xFA)

# ── Dimensions ────────────────────────────────────────────────
SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.8)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper Functions ──────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text="", font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(4), font_name="Calibri"):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_section_label(slide, text, top=Inches(0.5)):
    """Add the green uppercase section label."""
    add_text_box(slide, MARGIN, top, Inches(6), Inches(0.4),
                 text.upper(), font_size=11, color=ACCENT, bold=True)

def add_title(slide, text, top=Inches(0.9), font_size=36):
    """Add the main title."""
    add_text_box(slide, MARGIN, top, Inches(11), Inches(0.8),
                 text, font_size=font_size, color=WHITE, bold=True)

def add_subtitle(slide, text, top=Inches(1.7)):
    """Add subtitle text."""
    add_text_box(slide, MARGIN, top, Inches(10), Inches(0.6),
                 text, font_size=16, color=WHITE_60)

def stat_card(slide, left, top, width, height, number, label, description, number_color=ACCENT, border_color=None):
    """Create a stat card with big number + label + description."""
    card = add_rect(slide, left, top, width, height, BG_CARD, border_color or RGBColor(0x1A, 0x4A, 0x30))
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.6),
                 number, font_size=36, color=number_color, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.9), width - Inches(0.6), Inches(0.35),
                 label, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.3), width - Inches(0.6), Inches(height.inches - 1.6),
                 description, font_size=10, color=WHITE_60)

def info_card(slide, left, top, width, height, title_text, items, border_color=None):
    """Create an info card with title and bullet items."""
    card = add_rect(slide, left, top, width, height, BG_CARD, border_color or RGBColor(0x1A, 0x4A, 0x30))
    tb = add_text_box(slide, left + Inches(0.3), top + Inches(0.25), width - Inches(0.6), Inches(0.35),
                      title_text, font_size=14, color=WHITE, bold=True)
    y = top + Inches(0.7)
    for item in items:
        add_text_box(slide, left + Inches(0.3), y, width - Inches(0.6), Inches(0.3),
                     item, font_size=10, color=WHITE_60)
        y += Inches(0.28)

def accent_bar(slide, left, top, width, pct, label_left, label_right):
    """Create a progress bar with labels."""
    # Background bar
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.22))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x28)
    bar_bg.line.fill.background()
    bar_bg.adjustments[0] = 0.5
    # Fill bar
    fill_w = int(width * pct)
    if fill_w > 0:
        bar_fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, fill_w, Inches(0.22))
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = ACCENT
        bar_fill.line.fill.background()
        bar_fill.adjustments[0] = 0.5
    # Labels
    add_text_box(slide, left, top - Inches(0.25), Inches(3), Inches(0.25),
                 label_left, font_size=10, color=WHITE_60)
    add_text_box(slide, left + width - Inches(1), top - Inches(0.25), Inches(1), Inches(0.25),
                 label_right, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

add_text_box(slide, MARGIN, Inches(0.8), Inches(11), Inches(1.8),
             "The Green Revolution\nin Oaxaca", font_size=52, color=WHITE, bold=True)
# Color "Green Revolution" part — we'll add it as separate element
# Overlay the green text
add_text_box(slide, Inches(3.05), Inches(0.8), Inches(8), Inches(0.9),
             "Green Revolution", font_size=52, color=ACCENT, bold=True)

add_section_label(slide, "OUR VISION", top=Inches(2.8))
add_text_box(slide, MARGIN, Inches(3.2), Inches(10), Inches(0.7),
             "To become Mexico's first vertically integrated bamboo exporter —\nfrom plantation to the US construction market.",
             font_size=18, color=WHITE_80)

# Three feature boxes at bottom
features = [
    ("5-Year Head Start", "10 hectares fully planted since 2021"),
    ("High Barriers to Entry", "5-7 year growth cycle = permanent moat"),
    ("Strategic Location", "Hours from Mexico's newest Pacific port"),
]
for i, (title, desc) in enumerate(features):
    x = MARGIN + Inches(i * 3.8)
    card = add_rect(slide, x, Inches(4.8), Inches(3.5), Inches(1.5), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
    add_text_box(slide, x + Inches(0.3), Inches(5.0), Inches(2.9), Inches(0.35),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(5.4), Inches(2.9), Inches(0.6),
                 desc, font_size=10, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — WHY BAMBOO
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "WHY BAMBOO")
add_title(slide, "The World's Most Renewable\nConstruction Material")
add_subtitle(slide, "Not wood. Not plastic. A biological manufacturing platform that renews itself every year.")

stats = [
    ("91 cm", "Daily Growth", "Fastest-growing plant on Earth.\nHarvestable in 5 years,\nproduces for 50+."),
    ("2×", "Stronger Than Oak", "Guadua angustifolia: higher\ntensile strength than steel\nby weight. The 'vegetal steel.'"),
    ("12t", "CO₂ / Hectare / Year", "Bamboo sequesters 2× more\ncarbon than tropical forest.\nCarbon-negative material."),
    ("0", "Replanting Needed", "Harvest annually without\nreplanting. The root system\nregenerates indefinitely."),
]
for i, (num, label, desc) in enumerate(stats):
    x = MARGIN + Inches(i * 2.9)
    stat_card(slide, x, Inches(2.8), Inches(2.7), Inches(2.5), num, label, desc)

# Bottom insight box
insight = add_rect(slide, MARGIN, Inches(5.6), Inches(11.7), Inches(0.8), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, MARGIN + Inches(0.3), Inches(5.7), Inches(11.1), Inches(0.6),
             "Key insight: Bamboo is not agriculture — it's a perpetual asset that appreciates with time. Plant once, harvest forever.",
             font_size=12, color=WHITE_80, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "MARKET OPPORTUNITY")
add_title(slide, "A Massive, Growing Market")
add_subtitle(slide, "Global demand for engineered bamboo is growing 15-20% annually — supply can't keep up.")

markets = [
    ("$74.5B", "Global Bamboo Market", "2026 projection. CAGR 5.6%.\nConstruction is the fastest-\ngrowing segment.", "Total Addressable Market"),
    ("$10-12B", "US Engineered Bamboo", "Currently sourced almost\nexclusively from China and\nVietnam. US buyers actively\nseeking alternatives.", "Serviceable Market"),
    ("$0", "Mexican Structural Exports", "Zero Mexican companies\nexport domestically-grown\nstructural bamboo to the US.\nThe channel is completely\nunbuilt.", "White Space"),
]
for i, (num, label, desc, tag) in enumerate(markets):
    x = MARGIN + Inches(i * 3.9)
    stat_card(slide, x, Inches(2.6), Inches(3.6), Inches(3.5), num, label, desc)
    # Tag at bottom of card
    add_text_box(slide, x + Inches(0.3), Inches(5.4), Inches(3.0), Inches(0.3),
                 tag.upper(), font_size=9, color=ACCENT, bold=True)

add_text_box(slide, MARGIN, Inches(6.5), Inches(11.7), Inches(0.3),
             "Sources: Fortune Business Insights, UN Comtrade, Mexico DATAMEXICO trade data",
             font_size=8, color=WHITE_40, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — WHY NOW (INTEROCEANIC CORRIDOR)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "WHY NOW")
add_title(slide, "The Interoceanic Corridor")
add_subtitle(slide, "Mexico's $6-7B Panama Canal bypass puts us hours from the world's newest Pacific port.")

# Left card — Three Convergences
card_l = add_rect(slide, MARGIN, Inches(2.6), Inches(5.5), Inches(4.2), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, MARGIN + Inches(0.3), Inches(2.8), Inches(5), Inches(0.35),
             "Three Convergences", font_size=15, color=WHITE, bold=True)

convergences = [
    ("01", "USMCA Nearshoring Boom", "US-China decoupling creating massive demand for Mexican supply chains. Duty-free access to US market."),
    ("02", "Panama Canal in Crisis", "2023-24 drought cut transits 40%. Wait times 10-21+ days. Ships paying $1-4M for priority passage."),
    ("03", "Salina Cruz Port Expansion", "Mexico investing MXN 120B in the Interoceanic Corridor. New container terminal, 1.4M TEU capacity target."),
]
for i, (num, title, desc) in enumerate(convergences):
    y = Inches(3.3) + Inches(i * 1.15)
    add_text_box(slide, MARGIN + Inches(0.3), y, Inches(0.5), Inches(0.3),
                 num, font_size=12, color=ACCENT, bold=True)
    add_text_box(slide, MARGIN + Inches(0.8), y, Inches(4.5), Inches(0.3),
                 title, font_size=12, color=WHITE, bold=True)
    add_text_box(slide, MARGIN + Inches(0.8), y + Inches(0.3), Inches(4.5), Inches(0.5),
                 desc, font_size=9, color=WHITE_60)

# Right card — Logistics
card_r = add_rect(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(4.2), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.3), Inches(2.8), Inches(5), Inches(0.35),
             "Our Logistics Advantage", font_size=15, color=WHITE, bold=True)

routes = [
    ("Plantation → Salina Cruz", "~150 km · 3-5 hrs", ACCENT),
    ("Salina Cruz → US West Coast", "4-6 days", ACCENT),
    ("Rail to Gulf → Houston", "3-5 days", ACCENT),
    ("vs. Asia via Panama Canal", "2-3 weeks", RED_SOFT),
]
for i, (route, time, color) in enumerate(routes):
    y = Inches(3.4) + Inches(i * 0.55)
    add_text_box(slide, Inches(7.3), y, Inches(3.0), Inches(0.3),
                 route, font_size=10, color=WHITE_60)
    add_text_box(slide, Inches(10.3), y, Inches(2.0), Inches(0.3),
                 time, font_size=11, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

# Corridor info box
corr = add_rect(slide, Inches(7.3), Inches(5.8), Inches(4.9), Inches(0.8), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, Inches(7.5), Inches(5.85), Inches(4.5), Inches(0.3),
             "Corredor Interoceánico del Istmo de Tehuantepec", font_size=9, color=ACCENT, bold=True)
add_text_box(slide, Inches(7.5), Inches(6.15), Inches(4.5), Inches(0.3),
             "303 km rail · Salina Cruz (Pacific) ↔ Coatzacoalcos (Gulf) · 10 Industrial Poles with tax incentives",
             font_size=8, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — NO LOSS IN BAMBOO
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "ZERO WASTE MODEL")
add_title(slide, "No Loss in Bamboo")
add_subtitle(slide, "We target construction (highest value) — but every part of the plant has a market.")

# Primary revenue card (highlighted)
card1 = add_rect(slide, MARGIN, Inches(2.6), Inches(3.6), Inches(3.5), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, MARGIN + Inches(0.3), Inches(2.75), Inches(3.0), Inches(0.25),
             "PRIMARY REVENUE", font_size=9, color=ACCENT, bold=True)
add_text_box(slide, MARGIN + Inches(0.3), Inches(3.1), Inches(3.0), Inches(0.6),
             "400", font_size=48, color=ACCENT, bold=True)
add_text_box(slide, MARGIN + Inches(0.3), Inches(3.75), Inches(3.0), Inches(0.3),
             "MXN / pole", font_size=16, color=WHITE, bold=True)
add_text_box(slide, MARGIN + Inches(0.3), Inches(4.2), Inches(3.0), Inches(1.2),
             "Structural Premium — Treated Guadua poles for high-end construction. Hotels, villas, eco-luxury projects. Our core product and highest margin.",
             font_size=10, color=WHITE_60)
tag1 = add_rect(slide, MARGIN + Inches(0.3), Inches(5.5), Inches(3.0), Inches(0.35), RGBColor(0x15, 0x4A, 0x30), ACCENT)
add_text_box(slide, MARGIN + Inches(0.5), Inches(5.5), Inches(2.6), Inches(0.35),
             "Highest margin", font_size=10, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# Industrial bamboo card
card2 = add_rect(slide, Inches(4.7), Inches(2.6), Inches(3.6), Inches(3.5), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(5.0), Inches(2.75), Inches(3.0), Inches(0.25),
             "INDUSTRIAL BAMBOO", font_size=9, color=WHITE_40, bold=True)
add_text_box(slide, Inches(5.0), Inches(3.1), Inches(3.0), Inches(0.5),
             "Booming", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(5.0), Inches(3.6), Inches(3.0), Inches(0.3),
             "Market", font_size=16, color=WHITE_60, bold=True)
add_text_box(slide, Inches(5.0), Inches(4.1), Inches(3.0), Inches(1.4),
             "Straight sticks → furniture, panels, flooring, fencing. The industrial bamboo market is exploding globally. Every culm has value regardless of grade.",
             font_size=10, color=WHITE_60)
tag2 = add_rect(slide, Inches(5.0), Inches(5.5), Inches(3.0), Inches(0.35), RGBColor(0x1A, 0x3A, 0x28))
add_text_box(slide, Inches(5.2), Inches(5.5), Inches(2.6), Inches(0.35),
             "Growing 15-20% annually", font_size=10, color=WHITE_60, bold=True, alignment=PP_ALIGN.CENTER)

# Additional streams card
card3 = add_rect(slide, Inches(8.6), Inches(2.6), Inches(3.6), Inches(3.5), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(8.9), Inches(2.75), Inches(3.0), Inches(0.25),
             "ADDITIONAL STREAMS", font_size=9, color=WHITE_40, bold=True)
add_text_box(slide, Inches(8.9), Inches(3.1), Inches(3.0), Inches(0.5),
             "Every Part", font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(8.9), Inches(3.5), Inches(3.0), Inches(0.3),
             "Monetized", font_size=16, color=WHITE_60, bold=True)
streams = ["Carbon credits (~$21 USD/ha/yr)", "Nursery sales (seedlings)", "Charcoal & biochar", "Bamboo fiber & textiles", "Shoots (food industry)"]
for i, s in enumerate(streams):
    add_text_box(slide, Inches(8.9), Inches(4.1) + Inches(i * 0.28), Inches(3.0), Inches(0.28),
                 s, font_size=10, color=WHITE_60)

# Bottom insight
insight = add_rect(slide, MARGIN, Inches(6.3), Inches(11.7), Inches(0.7), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, MARGIN + Inches(0.3), Inches(6.35), Inches(11.1), Inches(0.55),
             "Downside protection: Even in our most conservative scenario, bamboo that doesn't meet structural grade still sells into industrial markets at strong margins.",
             font_size=11, color=WHITE_80, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "THE PROBLEM")
add_title(slide, "No Reliable Supply Chain Exists")
add_subtitle(slide, "US bamboo manufacturers depend on fragmented, uncertified, unreliable supply — mostly from Asia.")

problems = [
    ("01", "Fragmented Supply", "Mexican producers are small collectives (under 20 ha each). No single supplier can guarantee volume, consistency, or quality."),
    ("02", "No Vertical Integration", "Nobody controls growing + treatment + distribution. Growers sell raw, untreated poles. Builders source randomly."),
    ("03", "Zero Export Certification", "No Mexican producer has ICC-ES or ASTM certification. The US market door is closed without it."),
    ("04", "Asia Dependence = Risk", "US builders source from China/Vietnam. 30-45 day shipping, tariff uncertainty, ESG concerns. They want nearshore alternatives."),
]
for i, (num, title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = MARGIN + Inches(col * 6.0)
    y = Inches(2.7) + Inches(row * 2.1)
    card = add_rect(slide, x, y, Inches(5.6), Inches(1.8), BG_CARD, RGBColor(0x50, 0x2A, 0x2A))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.4),
                 num, font_size=22, color=RED_SOFT, bold=True)
    add_text_box(slide, x + Inches(0.9), y + Inches(0.2), Inches(4.4), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.9), y + Inches(0.65), Inches(4.4), Inches(0.9),
                 desc, font_size=10, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "COMPETITIVE LANDSCAPE")
add_title(slide, "Mexico's Bamboo Industry Today")
add_subtitle(slide, "Mostly collectives and micro-companies. Nobody controls the full chain.")

# Left — Current Players
add_text_box(slide, MARGIN, Inches(2.6), Inches(4), Inches(0.35),
             "Current Players", font_size=15, color=WHITE, bold=True)

competitors = [
    ("Guadua Selecto", "Chiapas · Grows + sells Guadua", "90 ha"),
    ("Bambuver", "Veracruz · Nursery + tourism + education", "Not export-focused"),
    ("Grupo Puebla Bambu", "Puebla · Cooperative cluster", "Fragmented"),
    ("Bambuterra / Qincha / Todo de Bambu", "Design studios · Don't grow bamboo", "Our customers"),
]
for i, (name, location, note) in enumerate(competitors):
    y = Inches(3.1) + Inches(i * 0.75)
    card = add_rect(slide, MARGIN, y, Inches(5.5), Inches(0.65), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
    add_text_box(slide, MARGIN + Inches(0.2), y + Inches(0.05), Inches(3.0), Inches(0.3),
                 name, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, MARGIN + Inches(0.2), y + Inches(0.3), Inches(3.0), Inches(0.25),
                 location, font_size=8, color=WHITE_40)
    nc = ACCENT if "customer" in note.lower() else WHITE_60
    add_text_box(slide, Inches(4.3), y + Inches(0.15), Inches(1.8), Inches(0.35),
                 note, font_size=9, color=nc, bold=("customer" in note.lower()), alignment=PP_ALIGN.RIGHT)

# Right — What They're Missing
add_text_box(slide, Inches(7.0), Inches(2.6), Inches(5), Inches(0.35),
             "What They're Missing", font_size=15, color=WHITE, bold=True)

card_r = add_rect(slide, Inches(7.0), Inches(3.1), Inches(5.5), Inches(3.0), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
missing = [
    "✗  Full value chain control (grow → process → sell)",
    "✗  Export certifications (ICC-ES / ASTM)",
    "✗  Scale for export (100+ ha minimum)",
    "✗  US market access or buyer relationships",
]
for i, item in enumerate(missing):
    add_text_box(slide, Inches(7.3), Inches(3.3) + Inches(i * 0.4), Inches(5.0), Inches(0.35),
                 item, font_size=11, color=WHITE_60 if i < 4 else ACCENT)

add_text_box(slide, Inches(7.3), Inches(4.95), Inches(5.0), Inches(0.35),
             "✓  Pacífico Bambu is building all of this", font_size=12, color=ACCENT, bold=True)

# Phase 2 note
phase2 = add_rect(slide, Inches(7.0), Inches(6.3), Inches(5.5), Inches(0.7), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.3), Inches(6.35), Inches(5.0), Inches(0.6),
             "Phase 2 strategy: Once we control the value chain, we sell seedlings to local farmers and buy back — creating a supply network we manage.",
             font_size=9, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "OUR SOLUTION")
add_title(slide, "Vertical Integration at Scale")
add_subtitle(slide, "Grow premium Guadua → treat on-site → sell direct. Control quality, pricing, and the entire value chain.")

# Value chain flow
steps = ["GROW", "HARVEST", "TREAT", "DRY", "SELL"]
descs = ["110 ha Guadua", "50 ha own + 60 leased", "Per-10ha pools", "Bodega network", "400+ MXN (low-end)"]
for i, (step, desc) in enumerate(zip(steps, descs)):
    x = MARGIN + Inches(i * 2.35)
    box = add_rect(slide, x, Inches(2.6), Inches(2.0), Inches(0.9), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
    add_text_box(slide, x + Inches(0.15), Inches(2.65), Inches(1.7), Inches(0.3),
                 step, font_size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.95), Inches(1.7), Inches(0.3),
                 desc, font_size=9, color=WHITE_60, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, x + Inches(2.0), Inches(2.8), Inches(0.35), Inches(0.4),
                     "→", font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Left — Planting Roadmap
card_l = add_rect(slide, MARGIN, Inches(3.8), Inches(5.5), Inches(3.2), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, MARGIN + Inches(0.3), Inches(3.95), Inches(5.0), Inches(0.3),
             "Planting Roadmap", font_size=14, color=WHITE, bold=True)

timeline = [
        ("2021 (done)", "10 ha", "Established"),
        ("2026", "0 ha", "Clear & prep 20 ha"),
        ("2027", "+15 ha", "First planting"),
        ("2028", "+15 ha", "Scale own land"),
        ("2029", "+10 ha", "Complete own (50 ha)"),
        ("2030-2031", "+60 ha", "LEASED from neighbors"),
        ("2032+", "110 ha", "Full operation"),
    ]
for i, (year, ha, note) in enumerate(timeline):
    y = Inches(4.4) + Inches(i * 0.4)
    c_yr = ACCENT if "2033" in year else WHITE_60
    c_ha = ACCENT if "190" in ha else WHITE
    add_text_box(slide, MARGIN + Inches(0.3), y, Inches(1.5), Inches(0.3), year, font_size=10, color=c_yr)
    add_text_box(slide, MARGIN + Inches(1.9), y, Inches(1.5), Inches(0.3), ha, font_size=10, color=c_ha, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, MARGIN + Inches(3.6), y, Inches(1.5), Inches(0.3), note, font_size=9, color=WHITE_40, alignment=PP_ALIGN.RIGHT)

# Right — Key Advantages
card_r = add_rect(slide, Inches(7.0), Inches(3.8), Inches(5.5), Inches(3.2), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.3), Inches(3.95), Inches(5.0), Inches(0.3),
             "Key Advantages", font_size=14, color=WHITE, bold=True)

advantages = [
    "40-year land lease + lease-and-replicate model with neighbors",
    "Own seedling propagation from 2029 — cost drops from 35 → 10 MXN (71% reduction)",
    "On-site treatment — boric acid/borax cold soak, 10 pools, 30 poles/batch × 7 days",
    "USMCA duty-free — Mexican-grown bamboo enters US without tariffs",
]
for i, adv in enumerate(advantages):
    y = Inches(4.4) + Inches(i * 0.65)
    add_text_box(slide, Inches(7.3), y, Inches(0.3), Inches(0.3), "✓", font_size=12, color=ACCENT)
    add_text_box(slide, Inches(7.7), y, Inches(4.5), Inches(0.55), adv, font_size=10, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — WHAT WE HAVE TODAY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "PROOF OF CONCEPT")
add_title(slide, "What We Have Today")
add_subtitle(slide, "We've already cleared the highest-risk phase. 5 years of biological establishment — complete.")

# Three hero stats
hero_stats = [
    ("10", "Hectares Planted", "Since April 2, 2021. 2,000 plants\nthriving. Guadua + Oldhamii + Gigante."),
    ("110", "Active Hectares Plan", "50 ha own + 60 ha leased.\nReplicable to 200+ ha."),
    ("5", "Years Established", "Past the biological risk phase.\nReady for first commercial harvest."),
]
for i, (num, label, desc) in enumerate(hero_stats):
    x = MARGIN + Inches(i * 3.9)
    card = add_rect(slide, x, Inches(2.6), Inches(3.6), Inches(2.1), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
    add_text_box(slide, x + Inches(0.3), Inches(2.75), Inches(3.0), Inches(0.7),
                 num, font_size=44, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.4), Inches(3.0), Inches(0.3),
                 label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.8), Inches(3.0), Inches(0.7),
                 desc, font_size=9, color=WHITE_60, alignment=PP_ALIGN.CENTER)

# Infrastructure
infra = add_rect(slide, MARGIN, Inches(5.0), Inches(5.5), Inches(2.0), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, MARGIN + Inches(0.3), Inches(5.15), Inches(5.0), Inches(0.3),
             "Infrastructure in Place", font_size=13, color=WHITE, bold=True)
infra_items = ["✓  Chevrolet 4×4 truck", "✓  Zapotal work yard + storage", "✓  Potable water system", "✓  Road access", "✓  2 bathrooms on-site"]
for i, item in enumerate(infra_items):
    col = i % 2
    row = i // 2
    add_text_box(slide, MARGIN + Inches(0.3) + Inches(col * 2.5), Inches(5.55) + Inches(row * 0.3),
                 Inches(2.3), Inches(0.3), item, font_size=9, color=WHITE_60)

# 2026 milestones
mile = add_rect(slide, Inches(7.0), Inches(5.0), Inches(5.5), Inches(2.0), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.3), Inches(5.15), Inches(5.0), Inches(0.3),
             "2026 Milestones (In Progress)", font_size=13, color=WHITE, bold=True)
miles = ["◉  Densify 10 ha (+3,500 seedlings)", "◉  Clear 20 ha for 2027 planting", "◉  First harvest: ~100 poles pilot", "◉  Build first drying bodega", "◉  Land title regularization"]
for i, m in enumerate(miles):
    add_text_box(slide, Inches(7.3), Inches(5.55) + Inches(i * 0.28),
                 Inches(5.0), Inches(0.28), m, font_size=9, color=YELLOW)



# ══════════════════════════════════════════════════════════════
# SLIDE 10 — REPLICATION / LEASE MODEL (Growth Lever)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "GROWTH LEVER")
add_title(slide, "The Lease-and-Replicate Model")
add_subtitle(slide, "We don't just operate land — we built a system that scales beyond the 110 ha. A blueprint for 200+ ha and beyond.")

# Three-phase model
phases = [
    ("PHASE 1", "OWN", "50 ha", "Pacífico Bambu's own land\n(within 190 ha lease)", ACCENT),
    ("PHASE 2", "LEASE", "+60 ha", "Neighbors' unused land\n10% revenue share after expenses", BLUE_SOFT),
    ("PHASE 3", "REPLICATE", "+100 ha", "Other locations across Oaxaca\nProven model, low capital", YELLOW),
]
for i, (label, action, ha_text, desc, color) in enumerate(phases):
    x = MARGIN + Inches(i * 4.0)
    card = add_rect(slide, x, Inches(2.6), Inches(3.8), Inches(2.4),
                    RGBColor(0x0F, 0x3A, 0x24) if i < 2 else BG_CARD,
                    color)
    add_text_box(slide, x + Inches(0.3), Inches(2.75), Inches(3.4), Inches(0.25),
                 label, font_size=10, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.05), Inches(3.4), Inches(0.45),
                 action, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.55), Inches(3.4), Inches(0.6),
                 ha_text, font_size=32, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(4.25), Inches(3.4), Inches(0.7),
                 desc, font_size=10, color=WHITE_60)

# Bottom — Why this works (4 pillars)
why_card = add_rect(slide, MARGIN, Inches(5.3), Inches(11.7), Inches(1.8), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, MARGIN + Inches(0.3), Inches(5.4), Inches(11.0), Inches(0.3),
             "Why The Replication Model Works", font_size=13, color=WHITE, bold=True)

pillars = [
    ("CAPITAL EFFICIENT", "No land purchase needed. Operations only — small CapEx per new site."),
    ("WIN-WIN ECONOMICS", "Neighbors monetize idle land. We get scale without massive equity raise."),
    ("RISK DISTRIBUTED", "Multiple locations = climate, pest, and operational risk diversified."),
    ("PROVEN PLAYBOOK", "By 2032 we'll have a tested system: workers, pools, bodegas, sales channel."),
]
for i, (title_p, desc_p) in enumerate(pillars):
    x = MARGIN + Inches(0.3) + Inches(i * 2.85)
    add_text_box(slide, x, Inches(5.8), Inches(2.7), Inches(0.25),
                 title_p, font_size=9, color=ACCENT, bold=True)
    add_text_box(slide, x, Inches(6.1), Inches(2.7), Inches(0.9),
                 desc_p, font_size=9, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — THE TEAM
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "THE TEAM")
add_title(slide, "Built to Execute")
add_subtitle(slide, "Local land expertise + professional bamboo construction + institutional growth capability.")

team = [
    ("J", "Juan", "Field Operations & Land Manager",
     "Lives in Candelaria. Manages all field operations, workers, planting, harvest, and land preparation. 5 years hands-on bamboo cultivation. Deep local knowledge and community relationships. The land is secured through Juan's 40-year lease.",
     ACCENT, True),
    ("D", "Diego", "Construction, Knowledge & Client Relations",
     "Most experienced team member. Professional bamboo constructor. Builds all structures and treatment infrastructure. Leads treatment protocols, quality control, and client relationships. Manages orders and delivery.",
     ACCENT, True),
    ("O", "Ofir", "Managing Partner & Head of Growth",
     "Strategy, finance, fundraising, investor relations, legal, marketing, branding, and certifications. Leads all business operations, financial modeling, and growth strategy. Executes the bridge-to-export roadmap.",
     ACCENT, True),
    ("R", "Rick", "Founding Capital Partner · First Angel",
     "First angel investor. Cash contributor with 20% equity stake. Skin in the game from day one — demonstrating conviction in the long-term thesis before external investors join.",
     WHITE_60, False),
]
for i, (initial, name, role, desc, color, highlighted) in enumerate(team):
    col = i % 2
    row = i // 2
    x = MARGIN + Inches(col * 6.0)
    y = Inches(2.7) + Inches(row * 2.3)
    border = ACCENT if highlighted else RGBColor(0x1A, 0x4A, 0x30)
    bg = RGBColor(0x0F, 0x3A, 0x24) if highlighted else BG_CARD
    card = add_rect(slide, x, y, Inches(5.6), Inches(2.0), bg, border)
    # Initial circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.25), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1A, 0x5A, 0x3A) if highlighted else RGBColor(0x1A, 0x3A, 0x28)
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.55), Inches(0.45),
                 initial, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Name and role
    add_text_box(slide, x + Inches(1.0), y + Inches(0.25), Inches(4.3), Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.55), Inches(4.3), Inches(0.25),
                 role, font_size=10, color=color)
    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(5.0), Inches(0.85),
                 desc, font_size=9, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — FINANCIAL PROJECTIONS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "FINANCIAL PROJECTIONS")
add_title(slide, "Path to Profitability")
add_subtitle(slide, "CONSERVATIVE BASE: 110 ha at maturity, 1,400 poles/ha mature yield, 400 MXN/pole. Options: 1K / 2K / 3K.")

# Bar chart area
chart_card = add_rect(slide, MARGIN, Inches(2.6), Inches(8.0), Inches(4.2), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, MARGIN + Inches(0.3), Inches(2.7), Inches(3), Inches(0.35),
             "Revenue Growth (MXN)", font_size=13, color=WHITE, bold=True)

bars = [
    ("2026", 0.034, 1), ("2027", 0.51, 2), ("2028", 0.68, 3), ("2029", 1.19, 5),
    ("2030", 1.39, 6), ("2032", 4.1, 18), ("2034", 15.0, 50), ("2038", 52.0, 100),
]
max_h = Inches(2.8)
for i, (year, revenue, pct) in enumerate(bars):
    x = MARGIN + Inches(0.5) + Inches(i * 0.9)
    bar_h = int(max_h * pct / 100)
    bar_y = Inches(6.1) - bar_h
    # Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, bar_y, Inches(0.6), bar_h)
    bar.fill.solid()
    opacity = 0.3 + (pct / 100) * 0.7
    g = int(0xD4 * opacity)
    b = int(0xA8 * opacity)
    bar.fill.fore_color.rgb = RGBColor(0x2D, g, b)
    bar.line.fill.background()
    bar.adjustments[0] = 0.15
    # Value label
    label = f"{revenue:.0f}M" if revenue >= 1 else f"{revenue*1000:.0f}K"
    is_highlight = year in ("2030", "2036")
    add_text_box(slide, x - Inches(0.1), bar_y - Inches(0.25), Inches(0.8), Inches(0.25),
                 label, font_size=8, color=ACCENT if is_highlight else WHITE_60, bold=is_highlight, alignment=PP_ALIGN.CENTER)
    # Year label
    add_text_box(slide, x - Inches(0.05), Inches(6.2), Inches(0.7), Inches(0.25),
                 year, font_size=8, color=ACCENT if is_highlight else WHITE_40, bold=is_highlight, alignment=PP_ALIGN.CENTER)

# Right side — milestones + cash-positive
mile_card = add_rect(slide, Inches(9.0), Inches(2.6), Inches(3.5), Inches(2.8), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(9.3), Inches(2.7), Inches(3.0), Inches(0.35),
             "Key Milestones", font_size=13, color=WHITE, bold=True)
milestones = [
    ("2026", "Pilot · 100 poles"),
    ("2027", "Ramp · 1,500 poles"),
    ("2031", "110 ha complete"),
    ("2033", "Cash-positive cum."),
    ("2038", "Full @ 1,400/ha"),
    ("—", "154K poles/year"),
]
for i, (yr, desc) in enumerate(milestones):
    y = Inches(3.15) + Inches(i * 0.35)
    c = ACCENT if yr == "2030" else WHITE_60
    add_text_box(slide, Inches(9.3), y, Inches(0.7), Inches(0.3), yr, font_size=10, color=c, bold=(yr=="2030"))
    add_text_box(slide, Inches(10.1), y, Inches(2.2), Inches(0.3), desc, font_size=9, color=c, bold=(yr=="2030"), alignment=PP_ALIGN.RIGHT)

# Cash-positive box
cp = add_rect(slide, Inches(9.0), Inches(5.6), Inches(3.5), Inches(1.2), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, Inches(9.3), Inches(5.65), Inches(3.0), Inches(0.2),
             "CASH-POSITIVE YEAR", font_size=9, color=ACCENT, bold=True)
add_text_box(slide, Inches(9.3), Inches(5.9), Inches(3.0), Inches(0.5),
             "2032", font_size=36, color=ACCENT, bold=True)
add_text_box(slide, Inches(9.3), Inches(6.4), Inches(3.0), Inches(0.3),
             "First profitable year.\nDividends from 2033+.", font_size=9, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — SENSITIVITY ANALYSIS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "SENSITIVITY ANALYSIS")
add_title(slide, "Three Scenarios")
add_subtitle(slide, "Even at conservative yields, the business is profitable. We model what happens if bamboo produces half of our base case.")

scenarios = [
    ("Conservative", "1,000", YELLOW, [("Revenue 2034", "10.6M MXN"), ("Revenue 2038", "37M MXN"), ("Break-even", "2034")], "Downside protection"),
    ("Base Case ⭐", "1,400", ACCENT, [("Revenue 2034", "15M MXN"), ("Revenue 2038", "52M MXN"), ("Break-even", "2033")], "OUR PROJECTION"),
    ("Optimistic", "2,000-3,000", BLUE_SOFT, [("Revenue 2034", "21-32M"), ("Revenue 2038", "73-110M"), ("Break-even", "2031-32")], "Possible upside"),
]
for i, (name, yield_val, color, metrics, tag) in enumerate(scenarios):
    x = MARGIN + Inches(i * 3.9)
    is_base = i == 1
    border = color if is_base else RGBColor(0x1A, 0x4A, 0x30)
    bg = RGBColor(0x0F, 0x3A, 0x24) if is_base else BG_CARD
    card = add_rect(slide, x, Inches(2.6), Inches(3.6), Inches(3.8), bg, border)
    # Dot + name
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(2.8), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, x + Inches(0.55), Inches(2.75), Inches(2.5), Inches(0.3),
                 name.upper(), font_size=10, color=color, bold=True)
    # Yield number
    add_text_box(slide, x + Inches(0.3), Inches(3.15), Inches(3.0), Inches(0.6),
                 yield_val, font_size=40, color=WHITE if not is_base else ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.7), Inches(3.0), Inches(0.3),
                 "poles/ha at maturity", font_size=10, color=WHITE_60)
    # Metrics
    for j, (metric_name, metric_val) in enumerate(metrics):
        y = Inches(4.3) + Inches(j * 0.4)
        add_text_box(slide, x + Inches(0.3), y, Inches(1.8), Inches(0.3), metric_name, font_size=10, color=WHITE_60)
        mc = ACCENT if "2032" in metric_val and is_base else WHITE
        add_text_box(slide, x + Inches(2.1), y, Inches(1.2), Inches(0.3), metric_val, font_size=10, color=mc, bold=True, alignment=PP_ALIGN.RIGHT)
    # Tag
    tag_box = add_rect(slide, x + Inches(0.3), Inches(5.8), Inches(3.0), Inches(0.35),
                       RGBColor(0x15, 0x2A, 0x20))
    add_text_box(slide, x + Inches(0.3), Inches(5.8), Inches(3.0), Inches(0.35),
                 tag, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, MARGIN, Inches(6.7), Inches(11.7), Inches(0.4),
             "CONSERVATIVE BASE: 1,400 poles/ha (not 2,000). Key assumptions: 400 MXN/pole · 15% waste · 110 ha (50 own + 60 leased). Only yield varies.",
             font_size=10, color=WHITE_60, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — THE ASK
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "INVESTMENT OPPORTUNITY")
add_title(slide, "The Ask")
add_subtitle(slide, "20% equity reserved for investors. First tranche at the highest-risk early valuation.")

# Left — Bridge Round
card_l = add_rect(slide, MARGIN, Inches(2.6), Inches(5.8), Inches(4.4), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, MARGIN + Inches(0.3), Inches(2.75), Inches(5.2), Inches(0.25),
             "BRIDGE ROUND — NOW OPEN", font_size=10, color=ACCENT, bold=True)
add_text_box(slide, MARGIN + Inches(0.3), Inches(3.1), Inches(5.2), Inches(0.7),
             "4,000,000", font_size=44, color=WHITE, bold=True)
add_text_box(slide, MARGIN + Inches(0.3), Inches(3.7), Inches(5.2), Inches(0.4),
             "MXN for 14% equity", font_size=18, color=WHITE_60)
add_text_box(slide, MARGIN + Inches(0.3), Inches(4.1), Inches(5.2), Inches(0.3),
             "Pre-money: ~24.5M MXN · Post-money: ~28.5M MXN", font_size=11, color=WHITE_40)

# Use of funds bars
add_text_box(slide, MARGIN + Inches(0.3), Inches(4.5), Inches(5.2), Inches(0.3),
             "Use of Funds", font_size=12, color=WHITE, bold=True)
funds = [
    ("4-year runway: ops to end-2029", 0.45),
    ("Land prep + planting 50 own ha", 0.25),
    ("Mgmt salaries (start 2027)", 0.20),
    ("Bodegas, pools, infrastructure", 0.10),
]
for i, (label, pct) in enumerate(funds):
    y = Inches(4.9) + Inches(i * 0.42)
    add_text_box(slide, MARGIN + Inches(0.3), y, Inches(3.5), Inches(0.2),
                 label, font_size=9, color=WHITE_60)
    add_text_box(slide, MARGIN + Inches(4.3), y, Inches(0.8), Inches(0.2),
                 f"{int(pct*100)}%", font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    # Bar
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     MARGIN + Inches(0.3), y + Inches(0.2), Inches(4.8), Inches(0.12))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x28)
    bar_bg.line.fill.background()
    bar_bg.adjustments[0] = 0.5
    bar_fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       MARGIN + Inches(0.3), y + Inches(0.2), Inches(4.8 * pct), Inches(0.12))
    bar_fill.fill.solid()
    opacity = 1.0 - (i * 0.2)
    bar_fill.fill.fore_color.rgb = RGBColor(0x2D, int(0xD4 * opacity), int(0xA8 * opacity))
    bar_fill.line.fill.background()
    bar_fill.adjustments[0] = 0.5

# Right — Equity + Valuation
eq_card = add_rect(slide, Inches(7.2), Inches(2.6), Inches(5.3), Inches(2.5), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.5), Inches(2.75), Inches(4.7), Inches(0.3),
             "Equity Allocation", font_size=13, color=WHITE, bold=True)

eq_items = [
    ("Founding partners (3)", "60%", Inches(2.4), ACCENT),
    ("Investors total (R1+R2)", "20%", Inches(0.8), BLUE_SOFT),
    ("Managing Partner (vested)", "15%", Inches(0.6), YELLOW),
    ("Strategic reserve", "5%", Inches(0.2), WHITE_40),
]
for i, (label, pct, bar_w, color) in enumerate(eq_items):
    y = Inches(3.2) + Inches(i * 0.5)
    add_text_box(slide, Inches(7.5), y, Inches(2.5), Inches(0.25), label, font_size=10, color=WHITE_60)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), y, bar_w, Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.5
    add_text_box(slide, Inches(10.2) + bar_w + Inches(0.1), y - Inches(0.02), Inches(0.6), Inches(0.25),
                 pct, font_size=10, color=WHITE, bold=True)

add_text_box(slide, Inches(7.5), Inches(4.6), Inches(4.7), Inches(0.4),
             "Bridge investor takes 14% — the largest piece. Funds 2026-2029 operations completely. Round 2 (Seed) closes 2028 at significantly higher valuation.",
             font_size=8, color=WHITE_40)

# Valuation box
val_card = add_rect(slide, Inches(7.2), Inches(5.3), Inches(5.3), Inches(1.7), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.5), Inches(5.45), Inches(4.7), Inches(0.3),
             "Why This Valuation Is Low", font_size=13, color=WHITE, bold=True)
add_text_box(slide, Inches(7.5), Inches(5.8), Inches(4.7), Inches(0.7),
             "You're entering before the risk phase is resolved. Once we prove production at scale and begin US export certification, the next round will reflect a fundamentally different risk profile — and a much higher valuation.",
             font_size=10, color=WHITE_60)
early = add_rect(slide, Inches(7.5), Inches(6.5), Inches(4.7), Inches(0.35), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, Inches(7.5), Inches(6.5), Inches(4.7), Inches(0.35),
             "Early entry = highest return potential", font_size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — NEXT RAISE / SEED
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "FUTURE GROWTH")
add_title(slide, "Round 2: The Seed")
add_subtitle(slide, "Smaller, lower-risk round in 2028 — funds the path to break-even at higher valuation.")

# Left — Strategic Investor Profile
card_l = add_rect(slide, MARGIN, Inches(2.6), Inches(5.8), Inches(4.2), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, MARGIN + Inches(0.3), Inches(2.75), Inches(5.2), Inches(0.25),
             "ROUND 2 — SEED (Q3 2028)", font_size=10, color=ACCENT, bold=True)
add_text_box(slide, MARGIN + Inches(0.3), Inches(3.2), Inches(5.2), Inches(0.8),
             "2,500,000 MXN for ~6% equity",
             font_size=22, color=WHITE, bold=True)

next_items = [
    "Funds 2030-2032 operations until break-even",
    "Lower risk: 2 years of production data + leased land model proven",
    "Pre-money: ~39M MXN — significant uplift from Round 1",
    "Profile: Agri-fund, Family Office, Impact Investor",
]
for i, item in enumerate(next_items):
    y = Inches(4.2) + Inches(i * 0.45)
    add_text_box(slide, MARGIN + Inches(0.3), y, Inches(0.3), Inches(0.3), "→", font_size=14, color=ACCENT)
    add_text_box(slide, MARGIN + Inches(0.7), y, Inches(4.8), Inches(0.4), item, font_size=10, color=WHITE_60)

# Right top — Why wait
why_card = add_rect(slide, Inches(7.2), Inches(2.6), Inches(5.3), Inches(2.3), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.5), Inches(2.75), Inches(4.7), Inches(0.3),
             "Why Wait for the Next Round?", font_size=13, color=WHITE, bold=True)
why_items = [
    ("01", "Round 1 funds 4 full years (2026-2029) — no rush for Round 2"),
    ("02", "By 2028 we have 1.02M MXN revenue + 25 ha producing — proven traction"),
    ("03", "Smaller raise = less dilution; founders keep more equity"),
]
for i, (num, text) in enumerate(why_items):
    y = Inches(3.2) + Inches(i * 0.55)
    add_text_box(slide, Inches(7.5), y, Inches(0.4), Inches(0.3), num, font_size=11, color=ACCENT, bold=True)
    add_text_box(slide, Inches(8.0), y, Inches(4.2), Inches(0.5), text, font_size=9, color=WHITE_60)

# Right bottom — Valuation trajectory
val_card = add_rect(slide, Inches(7.2), Inches(5.1), Inches(5.3), Inches(1.7), BG_CARD, RGBColor(0x1A, 0x4A, 0x30))
add_text_box(slide, Inches(7.5), Inches(5.25), Inches(4.7), Inches(0.3),
             "Valuation Trajectory", font_size=13, color=WHITE, bold=True)
add_text_box(slide, Inches(7.5), Inches(5.65), Inches(2.5), Inches(0.3),
             "Round 1 (2026)", font_size=10, color=WHITE_60)
add_text_box(slide, Inches(10.5), Inches(5.65), Inches(1.7), Inches(0.3),
             "28.5M MXN", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
add_text_box(slide, Inches(7.5), Inches(6.0), Inches(2.5), Inches(0.3),
             "Round 2 (2028)", font_size=10, color=WHITE_60)
add_text_box(slide, Inches(10.5), Inches(6.0), Inches(1.7), Inches(0.3),
             "~41.5M MXN", font_size=10, color=ACCENT, bold=True, alignment=PP_ALIGN.RIGHT)

uplift = add_rect(slide, Inches(7.5), Inches(6.4), Inches(4.7), Inches(0.3), RGBColor(0x0F, 0x3A, 0x24), ACCENT)
add_text_box(slide, Inches(7.5), Inches(6.4), Inches(4.7), Inches(0.3),
             "Round 1 enters at lowest valuation = highest upside potential", font_size=9, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 16 — MILESTONE ROADMAP
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_label(slide, "MILESTONE ROADMAP")
add_title(slide, "De-Risk, Then Scale")
add_subtitle(slide, "Every milestone is verifiable. Money deploys against proven gates, not promises.")

# Timeline bar
bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(2.6), Inches(11.7), Inches(0.08))
bar_bg.fill.solid()
bar_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x4A, 0x30)
bar_bg.line.fill.background()

# Timeline dots and labels
timeline_pts = [
    ("2026", 0.0), ("2027", 0.16), ("2028", 0.33),
    ("2030", 0.50), ("2032", 0.66), ("2034", 0.83), ("2038", 1.0),
]
for label, pos in timeline_pts:
    x = MARGIN + Inches(pos * 11.7)
    opacity = max(0.2, 1.0 - pos * 0.8)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.08), Inches(2.52), Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor(0x2D, int(0xD4 * opacity), int(0xA8 * opacity))
    dot.line.fill.background()
    add_text_box(slide, x - Inches(0.4), Inches(2.8), Inches(0.95), Inches(0.25),
                 label, font_size=8, color=WHITE_40, alignment=PP_ALIGN.CENTER)

# Milestone cards (2x3 grid)
milestones_data = [
    ("2026 — Foundation", "Deploy 4M Bridge. Densify 10 ha. Clear 20 ha for 2027. Build first drying bodega. First harvest: 100 poles. Land title regularization.", ACCENT, True),
    ("2027-2029 — Scale Own Land", "Plant 50 ha across 3 years (15+15+10). Build pools as needed. Revenue grows 510K → 1.7M MXN. Mgmt salaries start at 20K/month.", ACCENT, True),
    ("2028 — Round 2 Seed", "Close 2.5M Seed at higher valuation (~41.5M post). 2 years of production data. 25 ha already planted. Funds next 4 years.", ACCENT, True),
    ("2030-2031 — Lease Expansion", "Activate replication model: lease 60 ha from neighbors. 10% revenue share to landowners after expenses. 110 ha total active.", WHITE_60, False),
    ("2032 — Cash-Positive", "First profitable year. 4.1M MXN revenue. Vehicle purchases begin. Dividends planning starts.", WHITE_60, False),
    ("2033-2038 — Full Maturity", "All 110 ha producing. Revenue scales to 52M MXN (at 1,400/ha base). Big bodega 2034. ICC-ES 2037.", WHITE_40, False),
]
for i, (title, desc, color, highlighted) in enumerate(milestones_data):
    col = i % 2
    row = i // 2
    x = MARGIN + Inches(col * 6.0)
    y = Inches(3.2) + Inches(row * 1.35)
    border = ACCENT if highlighted else RGBColor(0x1A, 0x4A, 0x30)
    bg = RGBColor(0x0F, 0x3A, 0x24) if highlighted else BG_CARD
    card = add_rect(slide, x, y, Inches(5.6), Inches(1.2), bg, border)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(5.0), Inches(0.3),
                 title.upper(), font_size=10, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(5.0), Inches(0.65),
                 desc, font_size=9, color=WHITE_60)


# ══════════════════════════════════════════════════════════════
# SLIDE 17 — VISION / CLOSE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, MARGIN, Inches(1.5), Inches(11.7), Inches(1.0),
             "110 ha", font_size=64, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, MARGIN, Inches(2.5), Inches(11.7), Inches(0.5),
             "154,000 poles/year at full maturity · Replicable to 200+ ha", font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, MARGIN, Inches(3.0), Inches(11.7), Inches(0.5),
             "Mexico's first certified Guadua bamboo exporter", font_size=18, color=WHITE_60, alignment=PP_ALIGN.CENTER)

# Four pillars
pillars = [
    ("Hours from", "Salina Cruz Port"),
    ("Duty-free via", "USMCA"),
    ("40-year", "Land Lease"),
    ("Zero waste", "Model"),
]
for i, (sub, main) in enumerate(pillars):
    x = Inches(1.5) + Inches(i * 2.8)
    add_text_box(slide, x, Inches(4.0), Inches(2.2), Inches(0.25),
                 sub, font_size=9, color=WHITE_40, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(4.25), Inches(2.2), Inches(0.3),
                 main, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, x + Inches(2.3), Inches(4.1), Inches(0.3), Inches(0.4),
                     "|", font_size=16, color=RGBColor(0x1A, 0x4A, 0x30), alignment=PP_ALIGN.CENTER)

# Company name and tagline
add_text_box(slide, MARGIN, Inches(5.2), Inches(11.7), Inches(0.7),
             "Pacífico Bambu", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, MARGIN, Inches(5.9), Inches(11.7), Inches(0.4),
             "The bamboo supply chain from Mexico to the US — built right.", font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, MARGIN, Inches(6.6), Inches(11.7), Inches(0.3),
             "El Azulillo, Candelaria, Oaxaca, Mexico", font_size=11, color=WHITE_40, alignment=PP_ALIGN.CENTER)
add_text_box(slide, MARGIN, Inches(6.9), Inches(11.7), Inches(0.3),
             "contact@pacificobambu.com", font_size=11, color=WHITE_40, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "Pacifico_Bambu_Pitch_Deck_v3.pptx")
prs.save(output_path)
print(f"✓ Saved: {output_path}")
print(f"  {len(prs.slides)} slides generated")
